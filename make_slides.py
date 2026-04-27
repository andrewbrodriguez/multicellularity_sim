"""
Generates presentation.pptx — in-progress talk on the multicellularity sim.
Light theme, ~14 slides, video placeholders for Napari screen recordings.
Run:  python make_slides.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── light-theme palette ───────────────────────────────────────────────────────
BG          = RGBColor(0xff, 0xff, 0xff)    # paper white
PANEL       = RGBColor(0xf4, 0xf6, 0xfa)    # card background
PANEL_EDGE  = RGBColor(0xdc, 0xe0, 0xe8)
INK         = RGBColor(0x1a, 0x1d, 0x2e)    # primary text
MUTED       = RGBColor(0x5a, 0x5d, 0x70)    # secondary text
FAINT       = RGBColor(0x94, 0x97, 0xa6)    # dim labels
HAIRLINE    = RGBColor(0xe4, 0xe7, 0xed)

BLUE        = RGBColor(0x1f, 0x6f, 0xeb)    # cooperator
DEEP_BLUE   = RGBColor(0x14, 0x4e, 0xb5)
RED         = RGBColor(0xd1, 0x3a, 0x5e)    # defector
GREEN       = RGBColor(0x2e, 0x7d, 0x32)    # coop allele
AMBER       = RGBColor(0xe6, 0x7d, 0x00)    # accent
PURPLE      = RGBColor(0x7b, 0x3f, 0xbf)    # selection-pressure accent
PINK        = RGBColor(0xb5, 0x3a, 0x8c)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FIGS    = Path("figures")

# ── helpers ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill=PANEL, line=None, line_w=Pt(0), rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        l, t, w, h)
    if rounded:
        # reduce corner radius
        shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, text, l, t, w, h,
             size=Pt(18), bold=False, color=INK,
             align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top  = tf.margin_bottom = Pt(0)
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, l, t, w, h,
                body_size=Pt(16), sub_size=None,
                top_color=DEEP_BLUE, sub_color=MUTED):
    """items = list of (text, level, highlight_color or None)"""
    if sub_size is None:
        sub_size = Pt(body_size.pt - 2)
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top  = tf.margin_bottom = Pt(0)
    first = True
    for text, lvl, hi in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_before = Pt(5 if lvl == 0 else 2)
        p.space_after  = Pt(1)
        run = p.add_run()
        bullet = "▸ " if lvl == 0 else "    – "
        run.text = bullet + text
        run.font.size = body_size if lvl == 0 else sub_size
        run.font.bold = (lvl == 0)
        run.font.color.rgb = hi if hi else (top_color if lvl == 0 else sub_color)
    return box


def top_rule(slide, color=BLUE):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=color)


def section_label(slide, text, color=BLUE):
    add_text(slide, text.upper(),
             Inches(0.5), Inches(0.18), Inches(7), Inches(0.3),
             size=Pt(10), bold=True, color=color)


def slide_title(slide, title, subtitle=None):
    add_text(slide, title,
             Inches(0.5), Inches(0.38), Inches(12.3), Inches(0.7),
             size=Pt(28), bold=True, color=INK)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
                 size=Pt(14), color=MUTED)


def divider(slide, t, l=Inches(0.5), w=Inches(12.3), color=HAIRLINE):
    add_rect(slide, l, t, w, Pt(0.75), fill=color)


def page_number(slide, n, total):
    add_text(slide, f"{n} / {total}",
             Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3),
             size=Pt(9), color=FAINT, align=PP_ALIGN.RIGHT)


def embed_image(slide, path, l, t, w, h, caption=None):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), l, t, w, h)
    else:
        add_rect(slide, l, t, w, h, fill=PANEL, line=PANEL_EDGE, line_w=Pt(1))
        add_text(slide, f"[missing: {Path(path).name}]",
                 l + Inches(0.15), t + Inches(0.15),
                 w - Inches(0.3), h - Inches(0.3),
                 size=Pt(10), color=FAINT)
    if caption:
        add_text(slide, caption,
                 l, t + h + Inches(0.05), w, Inches(0.3),
                 size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER, italic=True)


def video_placeholder(slide, l, t, w, h, title, description, suggested_filename):
    """Frame with clear instructions for inserting a Napari screen recording."""
    # outer frame
    add_rect(slide, l, t, w, h, fill=RGBColor(0xfa, 0xfb, 0xfd),
             line=BLUE, line_w=Pt(2), rounded=True)
    # faint inner label
    add_text(slide, "VIDEO",
             l + Inches(0.15), t + Inches(0.1),
             w - Inches(0.3), Inches(0.3),
             size=Pt(10), bold=True, color=BLUE)
    # play triangle (rough, made from text)
    play_size = Inches(1.1)
    play_cx   = l + w / 2
    play_cy   = t + h / 2 - Inches(0.35)
    tri = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        play_cx - play_size / 2,
        play_cy - play_size / 2,
        play_size, play_size)
    tri.rotation = 90
    tri.fill.solid()
    tri.fill.fore_color.rgb = BLUE
    tri.line.fill.background()
    # title
    add_text(slide, title,
             l + Inches(0.15), t + h / 2 + Inches(0.35),
             w - Inches(0.3), Inches(0.4),
             size=Pt(15), bold=True, color=INK, align=PP_ALIGN.CENTER)
    # description
    add_text(slide, description,
             l + Inches(0.3), t + h / 2 + Inches(0.8),
             w - Inches(0.6), Inches(0.9),
             size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    # filename hint
    add_text(slide, f"drop in:  {suggested_filename}",
             l + Inches(0.3), t + h - Inches(0.4),
             w - Inches(0.6), Inches(0.3),
             size=Pt(9), color=FAINT, align=PP_ALIGN.CENTER)


def kv_stat(slide, l, t, w, value, label, value_color=BLUE):
    """Large-number stat card."""
    add_rect(slide, l, t, w, Inches(1.3), fill=PANEL,
             line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(slide, value,
             l, t + Inches(0.15), w, Inches(0.7),
             size=Pt(28), bold=True, color=value_color,
             align=PP_ALIGN.CENTER)
    add_text(slide, label,
             l, t + Inches(0.85), w, Inches(0.4),
             size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)


def chip(slide, l, t, text, color):
    """Small coloured chip (pill)."""
    w = Inches(0.05 * len(text) + 0.5)
    h = Inches(0.3)
    add_rect(slide, l, t, w, h, fill=color, rounded=True)
    add_text(slide, text,
             l, t + Inches(0.04), w, Inches(0.22),
             size=Pt(9), bold=True, color=RGBColor(0xff, 0xff, 0xff),
             align=PP_ALIGN.CENTER)
    return w


# ── SLIDE 1 — Title ──────────────────────────────────────────────────────────

def slide_title_page(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl)
    # coloured stripe
    add_rect(sl, 0, 0, Inches(0.3), SLIDE_H, fill=BLUE)
    # small label
    add_text(sl, "CS2212  ·  IN-PROGRESS REVIEW",
             Inches(0.7), Inches(1.0), Inches(8), Inches(0.3),
             size=Pt(11), bold=True, color=BLUE)
    # main title
    add_text(sl, "The Emergence of Multicellularity",
             Inches(0.7), Inches(1.6), Inches(11.5), Inches(1.2),
             size=Pt(42), bold=True, color=INK)
    add_text(sl, "A physics-based evolutionary simulation of the defector problem",
             Inches(0.7), Inches(2.9), Inches(11.5), Inches(0.6),
             size=Pt(20), color=MUTED)
    divider(sl, Inches(3.7), Inches(0.7), Inches(11.5))
    # author
    add_text(sl, "Andrew Rodriguez",
             Inches(0.7), Inches(3.9), Inches(8), Inches(0.4),
             size=Pt(18), bold=True, color=INK)
    add_text(sl, "April 2026   ·   Final project",
             Inches(0.7), Inches(4.35), Inches(8), Inches(0.3),
             size=Pt(13), color=MUTED)
    # chips
    x = Inches(0.7); y = Inches(5.3)
    for text, c in [("physics-based", BLUE), ("GPU-accelerated", AMBER),
                    ("evolutionary", GREEN), ("multilevel selection", PURPLE)]:
        w = chip(sl, x, y, text, c)
        x += w + Inches(0.15)
    page_number(sl, idx, total)


# ── SLIDE 2 — Motivation ──────────────────────────────────────────────────────

def slide_motivation(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, BLUE)
    section_label(sl, "Motivation", BLUE)
    slide_title(sl, "Why did cells ever start cooperating?",
                "Multicellularity is one of biology's major evolutionary transitions — and a paradox")

    # three columns of thought
    col_w = Inches(4.0); col_h = Inches(5.5); y = Inches(1.55); gap = Inches(0.2)
    xs = [Inches(0.5), Inches(0.5) + col_w + gap, Inches(0.5) + 2 * (col_w + gap)]

    # col 1 — fact
    add_rect(sl, xs[0], y, col_w, col_h, fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "01", xs[0] + Inches(0.25), y + Inches(0.2), Inches(0.6), Inches(0.4),
             size=Pt(12), bold=True, color=BLUE)
    add_text(sl, "Single cells dominated for ~2 billion years",
             xs[0] + Inches(0.25), y + Inches(0.55), col_w - Inches(0.5), Inches(1.0),
             size=Pt(16), bold=True, color=INK)
    add_text(sl, "Then, independently, lineages began forming persistent "
                 "cooperative clusters — animals, plants, fungi, algae. "
                 "Something changed the fitness landscape to favour coming together.",
             xs[0] + Inches(0.25), y + Inches(1.6), col_w - Inches(0.5), Inches(3.0),
             size=Pt(13), color=MUTED)

    # col 2 — tension
    add_rect(sl, xs[1], y, col_w, col_h, fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "02", xs[1] + Inches(0.25), y + Inches(0.2), Inches(0.6), Inches(0.4),
             size=Pt(12), bold=True, color=RED)
    add_text(sl, "Cooperation looks irrational for the individual",
             xs[1] + Inches(0.25), y + Inches(0.55), col_w - Inches(0.5), Inches(1.0),
             size=Pt(16), bold=True, color=INK)
    add_text(sl, "Paying a metabolic cost to benefit neighbours "
                 "is always undercut by a free-rider — a cell that receives "
                 "the shared benefit without paying in.",
             xs[1] + Inches(0.25), y + Inches(1.6), col_w - Inches(0.5), Inches(3.0),
             size=Pt(13), color=MUTED)

    # col 3 — question
    add_rect(sl, xs[2], y, col_w, col_h, fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "03", xs[2] + Inches(0.25), y + Inches(0.2), Inches(0.6), Inches(0.4),
             size=Pt(12), bold=True, color=AMBER)
    add_text(sl, "So what makes cooperation persist?",
             xs[2] + Inches(0.25), y + Inches(0.55), col_w - Inches(0.5), Inches(1.0),
             size=Pt(16), bold=True, color=INK)
    add_text(sl, "This project builds a simulation where we can tune "
                 "environmental conditions and directly observe when defectors "
                 "take over, when cooperators dominate, and what tips the balance.",
             xs[2] + Inches(0.25), y + Inches(1.6), col_w - Inches(0.5), Inches(3.0),
             size=Pt(13), color=MUTED)

    page_number(sl, idx, total)


# ── SLIDE 3 — The Defector Problem ────────────────────────────────────────────

def slide_defector(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, RED)
    section_label(sl, "The Core Problem", RED)
    slide_title(sl, "The Defector Dilemma",
                "Defectors join clusters, collect shared rewards, and contribute nothing")

    # left: text
    add_bullets(sl, [
        ("A cluster jointly computes a complex task", 0, DEEP_BLUE),
        ("Each contributing cell pays a compute cost (COOP_COST)", 1, None),
        ("The reward is split evenly across every member", 1, None),
        ("A defector exploits this split", 0, RED),
        ("Adhesion bit = 1  (can join clusters)", 1, None),
        ("Cooperator bit = 0  (does not contribute computation)", 1, None),
        ("Gets the full per-cell share of any reward the cluster earns", 1, None),
        ("Consequence inside a single cluster", 0, INK),
        ("Defectors have a strict per-tick fitness advantage", 1, None),
        ("Over time, mutation drives cooperators → defectors", 1, None),
        ("Cluster output collapses; the whole group starves", 1, None),
    ], l=Inches(0.5), t=Inches(1.55), w=Inches(6.3), h=Inches(5.5),
       body_size=Pt(15), top_color=INK)

    # right: panel showing the math
    add_rect(sl, Inches(7.0), Inches(1.55), Inches(5.8), Inches(5.5),
             fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "Fitness accounting",
             Inches(7.2), Inches(1.7), Inches(5.4), Inches(0.4),
             size=Pt(14), bold=True, color=INK)

    # equation-style layout
    eq_y = Inches(2.2)
    add_text(sl, "Cluster gross reward",
             Inches(7.2), eq_y, Inches(5.4), Inches(0.3),
             size=Pt(11), color=MUTED)
    add_text(sl, "R  =  Σ ( tile_reward × coop_scale )  −  drain × n_defectors",
             Inches(7.2), eq_y + Inches(0.3), Inches(5.4), Inches(0.4),
             size=Pt(12), bold=True, color=INK)

    add_text(sl, "Per-cell share",
             Inches(7.2), eq_y + Inches(0.85), Inches(5.4), Inches(0.3),
             size=Pt(11), color=MUTED)
    add_text(sl, "s  =  R  /  cluster_size",
             Inches(7.2), eq_y + Inches(1.15), Inches(5.4), Inches(0.4),
             size=Pt(12), bold=True, color=INK)

    add_text(sl, "Per-tick fitness change",
             Inches(7.2), eq_y + Inches(1.7), Inches(5.4), Inches(0.3),
             size=Pt(11), color=MUTED)
    add_text(sl, "cooperator :  +s  −  adhesion_cost  −  coop_cost",
             Inches(7.2), eq_y + Inches(2.0), Inches(5.4), Inches(0.35),
             size=Pt(12), bold=True, color=BLUE)
    add_text(sl, "defector      :  +s  −  adhesion_cost",
             Inches(7.2), eq_y + Inches(2.4), Inches(5.4), Inches(0.35),
             size=Pt(12), bold=True, color=RED)

    divider(sl, eq_y + Inches(2.95), Inches(7.2), Inches(5.4))

    add_text(sl, "The defector's edge is exactly coop_cost.\n"
                 "Cooperation only survives when group-level benefits overwhelm "
                 "the individual penalty — that's the regime we're hunting for.",
             Inches(7.2), eq_y + Inches(3.1), Inches(5.4), Inches(1.6),
             size=Pt(12), color=MUTED, italic=True)

    page_number(sl, idx, total)


# ── SLIDE 4 — Multilevel Selection ────────────────────────────────────────────

def slide_multilevel(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, PURPLE)
    section_label(sl, "Theoretical Framework", PURPLE)
    slide_title(sl, "Multilevel Selection",
                "Two opposing forces act on the same genome at the same time")

    # Two large opposing panels
    w = Inches(6.0); h = Inches(5.4); y = Inches(1.55)
    x_left  = Inches(0.5)
    x_right = Inches(0.5) + w + Inches(0.3)

    # individual selection (favours defectors)
    add_rect(sl, x_left, y, w, h, fill=PANEL, line=RED, line_w=Pt(1.5), rounded=True)
    add_text(sl, "WITHIN  CLUSTERS",
             x_left + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.3),
             size=Pt(10), bold=True, color=RED)
    add_text(sl, "Individual Selection  →  defectors win",
             x_left + Inches(0.3), y + Inches(0.6), w - Inches(0.6), Inches(0.6),
             size=Pt(18), bold=True, color=INK)
    add_bullets(sl, [
        ("Defectors skip COOP_COST every tick", 0, None),
        ("Mutation pressure biased toward defection", 0, None),
        ("    coop → def  at 2× base mutation rate", 1, None),
        ("    def → coop at 0.3× base rate", 1, None),
        ("Over generations, clusters drift toward defector-dominated", 0, None),
        ("Then produce no task reward → collapse", 0, None),
    ], l=x_left + Inches(0.3), t=y + Inches(1.35), w=w - Inches(0.6), h=Inches(3.6),
       body_size=Pt(14), top_color=INK, sub_color=MUTED)

    # group selection (favours cooperators)
    add_rect(sl, x_right, y, w, h, fill=PANEL, line=GREEN, line_w=Pt(1.5), rounded=True)
    add_text(sl, "BETWEEN  CLUSTERS",
             x_right + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.3),
             size=Pt(10), bold=True, color=GREEN)
    add_text(sl, "Group Selection  →  cooperators win",
             x_right + Inches(0.3), y + Inches(0.6), w - Inches(0.6), Inches(0.6),
             size=Pt(18), bold=True, color=INK)
    add_bullets(sl, [
        ("Cooperator-rich clusters earn higher total reward", 0, None),
        ("Whole-cluster replication — the group is the selection unit", 0, None),
        ("Functional clusters out-reproduce collapsing ones", 0, None),
        ("Spatial task specialisation amplifies the effect", 0, None),
        ("    a tile rewards (AND, XOR) — need BOTH operations", 1, None),
        ("    homogeneous clusters get zero on that tile", 1, None),
    ], l=x_right + Inches(0.3), t=y + Inches(1.35), w=w - Inches(0.6), h=Inches(3.6),
       body_size=Pt(14), top_color=INK, sub_color=MUTED)

    # bottom takeaway
    add_text(sl, "The key dial: how much does the environment reward group-level coordination "
                 "over individual efficiency?",
             Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
             size=Pt(13), color=PURPLE, align=PP_ALIGN.CENTER, italic=True)

    page_number(sl, idx, total)


# ── SLIDE 5 — Genome & Tasks ──────────────────────────────────────────────────

def slide_genome(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, GREEN)
    section_label(sl, "Model Design", GREEN)
    slide_title(sl, "Genome Encoding & Task Computation",
                "Every cell carries 8 bits; cooperation requires genomic complementarity")

    # left — genome strip
    add_rect(sl, Inches(0.5), Inches(1.5), Inches(6.3), Inches(2.4),
             fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "8-bit genome",
             Inches(0.7), Inches(1.65), Inches(6), Inches(0.35),
             size=Pt(13), bold=True, color=INK)

    bits = [
        ("0-1", "operation",  "AND / OR / XOR / NAND",  BLUE),
        ("2",   "adhesion",   "can join clusters",       AMBER),
        ("3",   "cooperator", "contributes computation", GREEN),
        ("4-7", "reserved",   "future traits",           FAINT),
    ]
    by = Inches(2.05)
    for i, (b, name, desc, col) in enumerate(bits):
        row_y = by + i * Inches(0.4)
        add_rect(sl, Inches(0.7), row_y, Inches(0.6), Inches(0.32),
                 fill=col, rounded=True)
        add_text(sl, b, Inches(0.7), row_y + Inches(0.04), Inches(0.6), Inches(0.25),
                 size=Pt(10), bold=True, color=RGBColor(0xff, 0xff, 0xff),
                 align=PP_ALIGN.CENTER)
        add_text(sl, name, Inches(1.4), row_y, Inches(1.6), Inches(0.32),
                 size=Pt(12), bold=True, color=INK)
        add_text(sl, desc, Inches(3.0), row_y, Inches(3.8), Inches(0.32),
                 size=Pt(11), color=MUTED)

    # phenotype rules box
    add_rect(sl, Inches(0.5), Inches(4.1), Inches(6.3), Inches(3.0),
             fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "Phenotype rules",
             Inches(0.7), Inches(4.25), Inches(6), Inches(0.35),
             size=Pt(13), bold=True, color=INK)
    add_bullets(sl, [
        ("Cooperator implies adhesion", 0, GREEN),
        ("You can't contribute to a group if you can't physically join one", 1, None),
        ("Cooperator = True  AND  Adhesion = True", 1, None),
        ("Defector = True  AND  NOT cooperator", 0, RED),
        ("Asymmetric mutation", 0, DEEP_BLUE),
        ("coop → def : 2× base  ·  def → coop : 0.3× base", 1, None),
        ("Reflects that breaking cooperation is genetically easier than building it", 1, None),
    ], l=Inches(0.7), t=Inches(4.65), w=Inches(6), h=Inches(2.4),
       body_size=Pt(13), sub_size=Pt(11), top_color=INK)

    # right — task computation
    add_rect(sl, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.6),
             fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "Multi-step tasks require DIFFERENT operations",
             Inches(7.2), Inches(1.65), Inches(5.5), Inches(0.4),
             size=Pt(14), bold=True, color=INK)
    add_text(sl, "Example — 2-step task  (AND, XOR)",
             Inches(7.2), Inches(2.1), Inches(5.5), Inches(0.3),
             size=Pt(11), color=MUTED)

    # pipeline arrows
    ty = Inches(2.55)
    # step boxes
    step_w = Inches(1.6); step_h = Inches(0.7)
    step1_x = Inches(7.2); step2_x = Inches(9.4); step3_x = Inches(11.1)
    add_rect(sl, step1_x, ty, step_w, step_h, fill=BLUE, rounded=True)
    add_text(sl, "cell₁\nAND (A, B)", step1_x, ty + Inches(0.05),
             step_w, step_h - Inches(0.1),
             size=Pt(11), bold=True, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)
    add_rect(sl, step2_x, ty, step_w, step_h, fill=AMBER, rounded=True)
    add_text(sl, "cell₂\nXOR (·, C)", step2_x, ty + Inches(0.05),
             step_w, step_h - Inches(0.1),
             size=Pt(11), bold=True, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)
    add_rect(sl, step3_x, ty, step_w, step_h, fill=GREEN, rounded=True)
    add_text(sl, "shared\nreward", step3_x, ty + Inches(0.05),
             step_w, step_h - Inches(0.1),
             size=Pt(11), bold=True, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)

    # arrow text
    add_text(sl, "→", step1_x + step_w, ty + Inches(0.15), Inches(0.4), Inches(0.5),
             size=Pt(20), color=MUTED, align=PP_ALIGN.CENTER)
    add_text(sl, "→", step2_x + step_w, ty + Inches(0.15), Inches(0.4), Inches(0.5),
             size=Pt(20), color=MUTED, align=PP_ALIGN.CENTER)

    add_bullets(sl, [
        ("A cluster of all-AND cooperators earns NOTHING", 0, RED),
        ("Even though they all \"cooperate\" genetically", 1, None),
        ("A cluster with one AND + one XOR unlocks the (AND, XOR) tile", 0, GREEN),
        ("Genome diversity is the requirement — not just the cooperator bit", 1, None),
        ("Tiles reward only multi-step tasks (indices 4-11)", 0, DEEP_BLUE),
        ("No reward for 1-step tasks → no incentive for solo work", 1, None),
    ], l=Inches(7.2), t=Inches(3.5), w=Inches(5.5), h=Inches(3.5),
       body_size=Pt(13), sub_size=Pt(11), top_color=INK)

    page_number(sl, idx, total)


# ── SLIDE 6 — Spatial Reward Landscape ────────────────────────────────────────

def slide_spatial(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, AMBER)
    section_label(sl, "Environment", AMBER)
    slide_title(sl, "Spatial Reward Landscape",
                "The 250µm × 250µm arena is tiled with task-specialised regions")

    add_bullets(sl, [
        ("Arena divided into spatial reward tiles", 0, DEEP_BLUE),
        ("Each tile carries a 12-slot reward vector (one per task)", 1, None),
        ("Drawn once at startup from a Dirichlet-Multinomial distribution", 1, None),
        ("Rewards are sparse — most tiles specialise in 1-2 tasks", 0, AMBER),
        ("Dirichlet α controls concentration:", 1, None),
        ("α = 0.01 → extreme specialisation (one tile rewards one task)", 1, None),
        ("α = 5.0  → uniform rewards everywhere (no spatial selection)", 1, None),
        ("Clusters must MOVE to where their operation is valued", 0, GREEN),
        ("A cluster that can cover many tasks thrives across regions", 1, None),
        ("A specialist cluster wins its tile but starves elsewhere", 1, None),
        ("Creates ecological niches — the foundation for diversity", 1, None),
    ], l=Inches(0.5), t=Inches(1.55), w=Inches(6.5), h=Inches(5.5),
       body_size=Pt(14), sub_size=Pt(12), top_color=INK)

    # right illustration: stylised 5x5 tile grid
    grid_x = Inches(7.3); grid_y = Inches(1.7)
    tile_w = Inches(1.0); gap = Inches(0.05); n = 5
    add_rect(sl, grid_x - Inches(0.1), grid_y - Inches(0.1),
             n * tile_w + (n - 1) * gap + Inches(0.2),
             n * tile_w + (n - 1) * gap + Inches(0.2),
             fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)

    # pseudo-random sparse tile colours (fixed for determinism)
    import random as _r
    _r.seed(7)
    palette = [BLUE, AMBER, GREEN, PURPLE, PINK, DEEP_BLUE, RED]
    for r in range(n):
        for c in range(n):
            tx = grid_x + c * (tile_w + gap)
            ty = grid_y + r * (tile_w + gap)
            if _r.random() < 0.35:
                col = _r.choice(palette)
                add_rect(sl, tx, ty, tile_w, tile_w, fill=col, rounded=False)
            else:
                add_rect(sl, tx, ty, tile_w, tile_w,
                         fill=RGBColor(0xed, 0xef, 0xf4), rounded=False)

    add_text(sl, "5 × 5 tiles of 50µm each.  Colour = dominant task.  Grey = low reward.",
             grid_x, grid_y + n * (tile_w + gap) + Inches(0.15),
             n * tile_w + (n - 1) * gap, Inches(0.3),
             size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    # tunable parameters callout
    cx = grid_x; cy = grid_y + n * (tile_w + gap) + Inches(0.65)
    cw = n * tile_w + (n - 1) * gap
    add_rect(sl, cx, cy, cw, Inches(1.2), fill=PANEL, line=AMBER, line_w=Pt(1), rounded=True)
    add_text(sl, "Tunable:",
             cx + Inches(0.2), cy + Inches(0.12), cw - Inches(0.4), Inches(0.3),
             size=Pt(11), bold=True, color=AMBER)
    add_text(sl, "task_alpha  ·  coop_reward_scale  ·  task_flip_period (MVG)",
             cx + Inches(0.2), cy + Inches(0.4), cw - Inches(0.4), Inches(0.7),
             size=Pt(13), bold=True, color=INK, align=PP_ALIGN.CENTER)

    page_number(sl, idx, total)


# ── SLIDE 7 — Physics Engine ──────────────────────────────────────────────────

def slide_physics(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, DEEP_BLUE)
    section_label(sl, "Implementation", DEEP_BLUE)
    slide_title(sl, "Physics Engine — Cell_Sim_2",
                "Gaussian pairwise force model with damped velocity integration")

    # left: force plot description + equations
    add_rect(sl, Inches(0.5), Inches(1.55), Inches(6.3), Inches(5.5),
             fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "Pairwise force  (§3.2)",
             Inches(0.7), Inches(1.7), Inches(6), Inches(0.35),
             size=Pt(13), bold=True, color=INK)
    add_text(sl, "F(d)  =  A(d)  −  50 · R(d)",
             Inches(0.7), Inches(2.1), Inches(6), Inches(0.5),
             size=Pt(22), bold=True, color=DEEP_BLUE, align=PP_ALIGN.CENTER)

    add_bullets(sl, [
        ("A(d) attractive Gaussian — centre 35 µm, σ = 5 µm", 0, GREEN),
        ("Only between cells in the SAME cluster", 1, None),
        ("No attraction across separate clusters — they stay independent", 1, None),
        ("R(d) repulsive Gaussian — centre −12.5 µm, σ = 10 µm", 0, RED),
        ("Applied to every pair — prevents overlap", 1, None),
        ("Net equilibrium ~24 µm centre-to-centre inside a cluster", 1, None),
    ], l=Inches(0.7), t=Inches(2.85), w=Inches(6), h=Inches(2.5),
       body_size=Pt(13), sub_size=Pt(11), top_color=INK)

    # velocity integration
    divider(sl, Inches(5.25), Inches(0.7), Inches(5.9))
    add_text(sl, "Integration step  (§2.3)",
             Inches(0.7), Inches(5.4), Inches(6), Inches(0.3),
             size=Pt(12), bold=True, color=INK)
    add_text(sl, "v  =  α · v_prev  +  M_i · F̂",
             Inches(0.7), Inches(5.8), Inches(6), Inches(0.4),
             size=Pt(16), bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(sl, "x  =  x + v  +  √(2D) · ξ     (Brownian diffusion)",
             Inches(0.7), Inches(6.2), Inches(6), Inches(0.4),
             size=Pt(14), bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(sl, "α = 0.12 damping · M_i = nearest-neighbour / 2 · max 4 µm/tick",
             Inches(0.7), Inches(6.65), Inches(6), Inches(0.3),
             size=Pt(10), color=FAINT, align=PP_ALIGN.CENTER, italic=True)

    # right: GPU stack
    add_rect(sl, Inches(7.0), Inches(1.55), Inches(5.8), Inches(5.5),
             fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "GPU acceleration via Taichi",
             Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.35),
             size=Pt(13), bold=True, color=AMBER)

    add_bullets(sl, [
        ("Spatial hash — 10 µm grid  (§2.1)", 0, DEEP_BLUE),
        ("25 × 25 buckets over the 250 µm arena", 1, None),
        ("Bucket assignment via atomic_add → lock-free parallel", 1, None),
        ("Force kernel scans 11×11 neighbourhood (110 µm reach)", 1, None),
        ("Two separate Taichi kernels per tick", 0, AMBER),
        ("_update_velocity:  sum pairwise forces, apply damping", 1, None),
        ("_update_position: integrate, add Brownian, reflect off walls", 1, None),
        ("Separation avoids read/write race conditions (paper §2.2)", 1, None),
        ("Hardware", 0, GREEN),
        ("Metal backend on Apple Silicon (falls back to CPU)", 1, None),
        ("~10× speedup over numpy for N > 500 cells", 1, None),
        ("Scales to 10,000+ cells in real-time", 1, None),
    ], l=Inches(7.2), t=Inches(2.1), w=Inches(5.5), h=Inches(5.0),
       body_size=Pt(13), sub_size=Pt(11), top_color=INK)

    page_number(sl, idx, total)


# ── SLIDE 8 — Live demo (video) ───────────────────────────────────────────────

def slide_demo_baseline(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, BLUE)
    section_label(sl, "Live Demo", BLUE)
    slide_title(sl, "What the simulation looks like",
                "Napari viewer — cells colour-coded by phenotype; scrub the time slider to replay")

    video_placeholder(sl,
        l=Inches(0.5), t=Inches(1.5), w=Inches(8.3), h=Inches(5.5),
        title="Baseline run — 500 ticks, default parameters",
        description="Lone cells scatter early (grey); adhesive cells bond into small clusters "
                    "(blue cooperators, red defectors); watch clusters drift across task tiles",
        suggested_filename="videos/baseline_run.mp4")

    # legend / side annotation
    add_rect(sl, Inches(9.0), Inches(1.5), Inches(3.8), Inches(5.5),
             fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "What to watch",
             Inches(9.2), Inches(1.65), Inches(3.5), Inches(0.35),
             size=Pt(13), bold=True, color=INK)

    # legend chips
    def legend_row(y_in, col, label, note):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(y_in), Inches(0.2), Inches(0.2))
        dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
        add_text(sl, label, Inches(9.5), Inches(y_in - 0.03), Inches(3.2), Inches(0.3),
                 size=Pt(11), bold=True, color=INK)
        add_text(sl, note, Inches(9.5), Inches(y_in + 0.22), Inches(3.2), Inches(0.4),
                 size=Pt(9), color=MUTED)

    # use "sl" not "slide" — fix local var
    def add_dot(y, col):
        dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.25), Inches(y), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()

    y = 2.15
    for col, label, note in [
        (FAINT, "lone cell",   "no cluster, drifts; earns base stipend"),
        (BLUE,  "cooperator",  "contributes compute, pays coop_cost"),
        (RED,   "defector",    "reaps share, pays nothing"),
        (AMBER, "cluster bond","visible link between adhered cells"),
    ]:
        add_dot(y, col)
        add_text(sl, label, Inches(9.5), Inches(y - 0.03), Inches(3.2), Inches(0.3),
                 size=Pt(11), bold=True, color=INK)
        add_text(sl, note, Inches(9.5), Inches(y + 0.22), Inches(3.2), Inches(0.4),
                 size=Pt(9), color=MUTED)
        y += 0.75

    page_number(sl, idx, total)


# ── SLIDE 9 — Cooperation Bonus Experiment ────────────────────────────────────

def slide_exp_coop_bonus(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, GREEN)
    section_label(sl, "Experiment 1", GREEN)
    slide_title(sl, "How strong does cooperation's reward need to be?",
                "Sweep: coop_reward_scale  ·  range  0.02 → 1.00  ·  3 seeds × 8 values")

    embed_image(sl, FIGS / "coop_bonus_coop_bonus_scatter.png",
                l=Inches(0.4), t=Inches(1.5), w=Inches(8.6), h=Inches(4.6),
                caption="coop_reward_scale vs final population composition, selection pressures, and cluster genome diversity")

    # findings panel on the right
    add_rect(sl, Inches(9.2), Inches(1.5), Inches(3.7), Inches(5.5),
             fill=PANEL, line=GREEN, line_w=Pt(1.2), rounded=True)
    add_text(sl, "Key finding",
             Inches(9.35), Inches(1.65), Inches(3.4), Inches(0.3),
             size=Pt(11), bold=True, color=GREEN)
    add_text(sl, "Sharp phase transition",
             Inches(9.35), Inches(1.95), Inches(3.4), Inches(0.5),
             size=Pt(17), bold=True, color=INK)
    add_bullets(sl, [
        ("Below ~0.10, cooperators collapse", 0, RED),
        ("coop_cost dominates the tiny shared reward", 1, None),
        ("Above ~0.15, cooperation stabilises", 0, GREEN),
        ("Group benefit exceeds individual cost", 1, None),
        ("Default (0.15) sits right at the knee", 0, AMBER),
        ("Small parameter shifts → big behavioural change", 1, None),
    ], l=Inches(9.35), t=Inches(2.55), w=Inches(3.4), h=Inches(4.3),
       body_size=Pt(12), sub_size=Pt(10), top_color=INK)

    # bottom key stats
    kv_stat(sl, Inches(0.4), Inches(6.3), Inches(2.1),  "≈ 0.12", "tipping point",  GREEN)
    kv_stat(sl, Inches(2.7), Inches(6.3), Inches(2.1),  "5% → 85%", "coop% jump",    BLUE)
    kv_stat(sl, Inches(5.0), Inches(6.3), Inches(2.1),  "24 runs", "trials (8 × 3)", DEEP_BLUE)
    kv_stat(sl, Inches(7.3), Inches(6.3), Inches(1.7),  "1500", "ticks / run",       AMBER)

    page_number(sl, idx, total)


# ── SLIDE 10 — Coop bonus videos side-by-side ─────────────────────────────────

def slide_coop_bonus_demo(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, GREEN)
    section_label(sl, "Live Demo", BLUE)
    slide_title(sl, "Below vs. above the cooperation threshold",
                "Same seed, two values of coop_reward_scale — the outcome diverges completely")

    video_placeholder(sl,
        l=Inches(0.4), t=Inches(1.55), w=Inches(6.3), h=Inches(5.0),
        title="coop_reward_scale = 0.05   —   cooperation fails",
        description="Clusters form but cooperators die out faster than they're born; "
                    "defectors dominate; clusters collapse into lone cells",
        suggested_filename="videos/coop_scale_low.mp4")

    video_placeholder(sl,
        l=Inches(6.9), t=Inches(1.55), w=Inches(6.0), h=Inches(5.0),
        title="coop_reward_scale = 0.40   —   cooperation thrives",
        description="Multi-operation clusters stabilise; cooperators fill the arena; "
                    "defectors stay rare and localised",
        suggested_filename="videos/coop_scale_high.mp4")

    # takeaway strip
    add_rect(sl, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.5),
             fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "Same physics. Same mutation rate. Same seed. Only the reward multiplier changed — "
                 "and a whole evolutionary outcome flipped.",
             Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.35),
             size=Pt(12), color=INK, italic=True, align=PP_ALIGN.CENTER)

    page_number(sl, idx, total)


# ── SLIDE 11 — Task distribution experiment ───────────────────────────────────

def slide_exp_task_dist(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, AMBER)
    section_label(sl, "Experiment 2", AMBER)
    slide_title(sl, "Does spatial task specialisation select for genome diversity?",
                "Sweep: task_alpha (Dirichlet concentration)  ·  0.01 (harsh) → 5.0 (uniform)")

    embed_image(sl, FIGS / "task_distribution_task_distribution.png",
                l=Inches(0.4), t=Inches(1.5), w=Inches(8.6), h=Inches(4.6),
                caption="task_alpha vs population mix, selection pressures, cluster op-diversity, and Shannon entropy")

    add_rect(sl, Inches(9.2), Inches(1.5), Inches(3.7), Inches(5.5),
             fill=PANEL, line=AMBER, line_w=Pt(1.2), rounded=True)
    add_text(sl, "Key finding",
             Inches(9.35), Inches(1.65), Inches(3.4), Inches(0.3),
             size=Pt(11), bold=True, color=AMBER)
    add_text(sl, "Specialisation breeds diversity",
             Inches(9.35), Inches(1.95), Inches(3.4), Inches(0.5),
             size=Pt(17), bold=True, color=INK)
    add_bullets(sl, [
        ("Low α (sparse tiles)", 0, AMBER),
        ("Clusters need ≥3 distinct ops to harvest local tile", 1, None),
        ("Strong group selection pressure", 1, None),
        ("High α (uniform tiles)", 0, RED),
        ("Any cluster earns on any tile", 1, None),
        ("No incentive for complementarity → defectors thrive", 1, None),
        ("Cluster genome entropy tracks α almost linearly", 0, DEEP_BLUE),
    ], l=Inches(9.35), t=Inches(2.55), w=Inches(3.4), h=Inches(4.3),
       body_size=Pt(12), sub_size=Pt(10), top_color=INK)

    # bottom stats
    kv_stat(sl, Inches(0.4), Inches(6.3), Inches(2.1),  "α ≤ 0.05", "pro-coop regime", AMBER)
    kv_stat(sl, Inches(2.7), Inches(6.3), Inches(2.1),  "α ≥ 1.0",  "pro-defect regime", RED)
    kv_stat(sl, Inches(5.0), Inches(6.3), Inches(2.1),  "3.2",       "ops/cluster at α=0.01", GREEN)
    kv_stat(sl, Inches(7.3), Inches(6.3), Inches(1.7),  "1.1",       "ops/cluster at α=5", MUTED)

    page_number(sl, idx, total)


# ── SLIDE 12 — Task distribution videos ───────────────────────────────────────

def slide_task_dist_demo(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, AMBER)
    section_label(sl, "Live Demo", BLUE)
    slide_title(sl, "Harsh specialisation vs. flat rewards",
                "Watch how the spatial reward structure shapes cluster composition")

    video_placeholder(sl,
        l=Inches(0.4), t=Inches(1.55), w=Inches(6.3), h=Inches(5.0),
        title="task_alpha = 0.01  —  harsh specialisation",
        description="Each tile rewards one task; clusters visibly shuttle to their home regions, "
                    "and internal genome diversity is high",
        suggested_filename="videos/task_alpha_low.mp4")

    video_placeholder(sl,
        l=Inches(6.9), t=Inches(1.55), w=Inches(6.0), h=Inches(5.0),
        title="task_alpha = 5.0  —  uniform rewards",
        description="No tile cares which task you do; clusters are generic and small; "
                    "defector invasions are common",
        suggested_filename="videos/task_alpha_high.mp4")

    add_rect(sl, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.5),
             fill=PANEL, line=PANEL_EDGE, line_w=Pt(0.5), rounded=True)
    add_text(sl, "Spatial structure isn't cosmetic — it's the substrate that lets multicellularity pay off.",
             Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.35),
             size=Pt(12), color=INK, italic=True, align=PP_ALIGN.CENTER)

    page_number(sl, idx, total)


# ── SLIDE 13 — What I'm focusing on next ──────────────────────────────────────

def slide_focus_next(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, PURPLE)
    section_label(sl, "What's next — core focus", PURPLE)
    slide_title(sl, "Where I'm taking this",
                "The infrastructure is complete; these are the science questions I'm actively pursuing")

    # three focus cards
    card_w = Inches(4.1); card_h = Inches(5.4); y = Inches(1.55); gap = Inches(0.15)
    xs = [Inches(0.5), Inches(0.5) + card_w + gap, Inches(0.5) + 2 * (card_w + gap)]

    def focus_card(x, tag, color, title, bullets):
        add_rect(sl, x, y, card_w, card_h, fill=PANEL,
                 line=color, line_w=Pt(1.5), rounded=True)
        add_rect(sl, x, y, Inches(0.6), Inches(0.35), fill=color, rounded=True)
        add_text(sl, tag, x + Inches(0.05), y + Inches(0.04),
                 Inches(0.55), Inches(0.3), size=Pt(10), bold=True,
                 color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)
        add_text(sl, title, x + Inches(0.25), y + Inches(0.55),
                 card_w - Inches(0.5), Inches(0.8),
                 size=Pt(16), bold=True, color=INK)
        add_bullets(sl, bullets,
                    l=x + Inches(0.25), t=y + Inches(1.3),
                    w=card_w - Inches(0.5), h=card_h - Inches(1.5),
                    body_size=Pt(12), sub_size=Pt(10),
                    top_color=INK, sub_color=MUTED)

    focus_card(xs[0], "A", BLUE, "Modularly-Varying Goals",
        [("The hypothesis", 0, DEEP_BLUE),
         ("If the rewarding task structure flips periodically, selection should favour "
          "genomes that can rapidly reconfigure", 1, None),
         ("What I'll measure", 0, DEEP_BLUE),
         ("Cluster op-diversity vs. flip period", 1, None),
         ("Post-flip recovery time for coop%", 1, None),
         ("Does modularity in genome space correspond to resilience?", 1, None)])

    focus_card(xs[1], "B", RED, "Defector invasion dynamics",
        [("The hypothesis", 0, DEEP_BLUE),
         ("A newly-formed cooperator cluster is metastable — "
          "mutation eventually produces a defector inside it", 1, None),
         ("What I'll measure", 0, DEEP_BLUE),
         ("Time from cluster birth to first defector", 1, None),
         ("Time from first defector to collapse", 1, None),
         ("How kin-selection (genome-similarity bonding) changes both", 1, None)])

    focus_card(xs[2], "C", GREEN, "Joint parameter phase diagram",
        [("The hypothesis", 0, DEEP_BLUE),
         ("The two experiments so far varied one axis each — "
          "the real structure lives in the 2D (α × coop_scale) plane", 1, None),
         ("What I'll measure", 0, DEEP_BLUE),
         ("Grid sweep over both parameters", 1, None),
         ("Map regimes: defector, cooperative, extinction, oscillating", 1, None),
         ("Identify the phase boundary analytically if possible", 1, None)])

    page_number(sl, idx, total)


# ── SLIDE 14 — Broader future directions ─────────────────────────────────────

def slide_future(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl); top_rule(sl, MUTED)
    section_label(sl, "Beyond the core focus", MUTED)
    slide_title(sl, "Longer-horizon directions",
                "Questions I'd love to explore if time permits")

    add_bullets(sl, [
        ("Emergent division of labour", 0, GREEN),
        ("Do successful clusters converge on complementary operation pairs (AND+XOR, OR+NAND)?", 1, None),
        ("Can we detect proto-specialisation from genome phylogenies across replicate runs?", 1, None),
        ("Spatial ecology", 0, BLUE),
        ("Do cooperator clusters form stable territories against lone-defector invaders?", 1, None),
        ("Do ecological niches emerge naturally from Dirichlet-sparse rewards?", 1, None),
        ("Longer evolutionary horizons", 0, AMBER),
        ("Current runs: 1,000–5,000 ticks.  Target: 50,000+ to observe cycling", 1, None),
        ("Do cooperation and defection oscillate (Red Queen dynamics)?", 1, None),
        ("Further genome traits", 0, PURPLE),
        ("Bits 4–7 are reserved — candidates: signalling, apoptosis, kin-recognition", 1, None),
        ("Does signalling between cooperators break the defector's cover?", 1, None),
        ("Better analytics", 0, DEEP_BLUE),
        ("Lineage tracking across runs — build phylogenetic trees from cell_id ancestry", 1, None),
        ("Automatic regime classification from time-series shape", 1, None),
    ], l=Inches(0.5), t=Inches(1.55), w=Inches(12.3), h=Inches(5.7),
       body_size=Pt(14), sub_size=Pt(12), top_color=INK)

    page_number(sl, idx, total)


# ── SLIDE 15 — Summary ────────────────────────────────────────────────────────

def slide_closing(prs, idx, total):
    sl = blank_slide(prs); fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.3), SLIDE_H, fill=BLUE)

    add_text(sl, "Summary",
             Inches(0.7), Inches(0.7), Inches(12), Inches(0.6),
             size=Pt(32), bold=True, color=INK)

    add_bullets(sl, [
        ("Built a GPU-accelerated, physics-based evolutionary simulation of multicellularity", 0, INK),
        ("Cell_Sim_2 Gaussian force model · Taichi Metal backend · O(N) spatial hash", 1, None),
        ("Task computation requires DIFFERENT operations from DIFFERENT cells — "
         "homogeneous clusters earn nothing", 0, GREEN),
        ("Cooperation is enforced at the phenotype level, not just named", 1, None),
        ("Two experiments completed:", 0, BLUE),
        ("Cooperation emerges past a sharp reward-scale threshold (~0.12)", 1, None),
        ("Spatial task specialisation selects for cluster genome diversity", 1, None),
        ("Next up: MVG task-flipping, defector invasion timing, 2D phase diagram "
         "over (α × coop_scale)", 0, PURPLE),
        ("Infrastructure ready: batch sweep script, CSV export, publication-quality figure tooling", 1, None),
    ], l=Inches(0.7), t=Inches(1.6), w=Inches(12), h=Inches(4.6),
       body_size=Pt(16), sub_size=Pt(13), top_color=INK)

    divider(sl, Inches(6.3), Inches(0.7), Inches(12))
    add_text(sl, "Thank you  —  questions welcome",
             Inches(0.7), Inches(6.5), Inches(8), Inches(0.5),
             size=Pt(18), bold=True, color=INK)
    add_text(sl, "Andrew Rodriguez   ·   andrewbrodriguez@gmail.com",
             Inches(0.7), Inches(7.0), Inches(8), Inches(0.35),
             size=Pt(11), color=MUTED)

    page_number(sl, idx, total)


# ── assemble ──────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    builders = [
        slide_title_page,
        slide_motivation,
        slide_defector,
        slide_multilevel,
        slide_genome,
        slide_spatial,
        slide_physics,
        slide_demo_baseline,
        slide_exp_coop_bonus,
        slide_coop_bonus_demo,
        slide_exp_task_dist,
        slide_task_dist_demo,
        slide_focus_next,
        slide_future,
        slide_closing,
    ]
    total = len(builders)
    for i, b in enumerate(builders, 1):
        b(prs, i, total)

    out = Path("presentation.pptx")
    prs.save(out)
    print(f"Saved {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
